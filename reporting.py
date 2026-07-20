from __future__ import annotations

from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def create_word_report(path: Path, results: pd.DataFrame, match_summary: pd.DataFrame, figures: list[Path]) -> None:
    doc=Document(); doc.add_heading('Hospital Quality and Insurance Base Rate Analysis',0)
    doc.add_paragraph('Objective: Determine whether hospitals with higher quality recognition receive higher insurance base rates for inpatient spine surgery DRGs.')
    doc.add_heading('Data integration',level=1)
    for _,row in match_summary.iterrows(): doc.add_paragraph(f"{row['Metric']}: {row['Value']}")
    doc.add_heading('Statistical results',level=1)
    if results.empty: doc.add_paragraph('No statistical test met the minimum data requirements.')
    else:
        table=doc.add_table(rows=1,cols=5); hdr=table.rows[0].cells
        for i,h in enumerate(['Analysis','Estimate','P value','N','Interpretation']): hdr[i].text=h
        for _,r in results.iterrows():
            cells=table.add_row().cells
            cells[0].text=str(r.get('Analysis','')); cells[1].text=f"{r.get('Estimate',float('nan')):.4g}" if pd.notna(r.get('Estimate')) else ''
            cells[2].text=f"{r.get('P_Value',float('nan')):.4g}" if pd.notna(r.get('P_Value')) else ''; cells[3].text=str(r.get('N','')); cells[4].text=str(r.get('Interpretation',''))
    for fig in figures:
        doc.add_heading(fig.stem.replace('_',' ').title(),level=2); doc.add_picture(str(fig),width=Inches(6.2))
    doc.save(path)


def create_powerpoint(path: Path, results: pd.DataFrame, match_summary: pd.DataFrame, figures: list[Path]) -> None:
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[0]); slide.shapes.title.text='Hospital Quality and Insurance Base Rates'; slide.placeholders[1].text='Python healthcare analytics project'
    slide=prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text='Project Aim'; slide.placeholders[1].text='Determine whether Blue Distinction status and CMS Overall Hospital Star Ratings are associated with higher hospital base rates for inpatient spine surgery DRGs.'
    slide=prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text='Data Matching Summary'; slide.placeholders[1].text="\n".join(f"{r['Metric']}: {r['Value']}" for _,r in match_summary.iterrows())
    slide=prs.slides.add_slide(prs.slide_layouts[1]); slide.shapes.title.text='Key Statistical Results'; slide.placeholders[1].text="\n".join(f"{r['Analysis']}: estimate={r['Estimate']:.3g}, p={r['P_Value']:.3g}" for _,r in results.iterrows()) if not results.empty else 'No tests met minimum data requirements.'
    for fig in figures[:3]:
        slide=prs.slides.add_slide(prs.slide_layouts[5]); slide.shapes.title.text=fig.stem.replace('_',' ').title(); slide.shapes.add_picture(str(fig),Inches(1.1),Inches(1.5),width=Inches(7.8))
    prs.save(path)
